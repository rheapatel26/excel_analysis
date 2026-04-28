"""
ppt_generator.py – Generates an editable PowerPoint presentation from analysis data.
Recreates the Corporate Health Insurance Review template layout with dynamic data.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io
import math

# ─── Color palette (matching the template) ─────────────────────────────────
NAVY = RGBColor(0x11, 0x18, 0x27)
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x3B)
LIGHT_GREEN = RGBColor(0x34, 0xD3, 0x99)
COPPER = RGBColor(0xB8, 0x73, 0x33)
LIGHT_COPPER = RGBColor(0xD4, 0x9B, 0x6A)
BG_GRAY = RGBColor(0xF0, 0xED, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MED = RGBColor(0x64, 0x74, 0x8B)
TEXT_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
DANGER_RED = RGBColor(0xF8, 0x71, 0x71)
ACCENT_BLUE = RGBColor(0x6C, 0x8F, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
TABLE_HEADER_BG = RGBColor(0x1E, 0x29, 0x3B)
TABLE_ALT_BG = RGBColor(0xF8, 0xFA, 0xFC)
SLIDE_BG = BG_GRAY

# Slide dimensions (widescreen 16:9 standard - 13.333 x 7.5 inches)
SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)


def _fmt_inr(amount):
    """Format amount in Indian Rupee notation (lakhs/crores)."""
    if amount >= 1e7:
        return f"₹{amount/1e7:.2f} Cr"
    elif amount >= 1e5:
        return f"₹{amount/1e5:.2f} L"
    else:
        return f"₹{amount:,.0f}"


def _fmt_inr_full(amount):
    """Format full Indian number with commas."""
    if amount < 0:
        return f"-₹{abs(amount):,.0f}"
    return f"₹{amount:,.0f}"


BG_GRAY = RGBColor(0xF8, 0xFA, 0xFC)  # Lighter, premium slate gray
SLIDE_BG = BG_GRAY

def _set_slide_bg(slide, color=SLIDE_BG):
    """Set solid background fill for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color
    
    # Premium branding accent line across the top of every slide
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), fill_color=NAVY)
    _add_shape_rect(slide, Inches(0), Inches(0.08), SLIDE_WIDTH, Inches(0.02), fill_color=ACCENT_BLUE)


def _add_textbox(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def _add_shape_rect(slide, left, top, width, height, fill_color=WHITE,
                    border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color=WHITE):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_kpi_card(slide, left, top, width, height, icon, value, label,
                  accent_color=ACCENT_BLUE, sub_text=None):
    """Add a KPI card with accent bar, value, and label."""
    # Card background
    card = _add_rounded_rect(slide, left, top, width, height, fill_color=WHITE)

    # Top accent bar
    _add_shape_rect(slide, left + Inches(0.05), top + Inches(0.05),
                    width - Inches(0.1), Inches(0.04), fill_color=accent_color)

    # Icon
    _add_textbox(slide, left + Inches(0.2), top + Inches(0.2),
                 Inches(0.5), Inches(0.5), icon, font_size=24)

    # Value
    _add_textbox(slide, left + Inches(0.2), top + Inches(0.65),
                 width - Inches(0.4), Inches(0.5), str(value),
                 font_size=28, bold=True, color=accent_color)

    # Label
    _add_textbox(slide, left + Inches(0.2), top + Inches(1.15),
                 width - Inches(0.4), Inches(0.35), label,
                 font_size=10, color=TEXT_MED)

    # Sub-text
    if sub_text:
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.45),
                     width - Inches(0.4), Inches(0.4), sub_text,
                     font_size=9, color=TEXT_LIGHT)


def _add_table(slide, left, top, width, headers, rows, col_widths=None):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                          Inches(0.35 * n_rows))
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            # Alternate row coloring
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = TEXT_DARK
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER

    return table_shape


def _add_insight_bar(slide, left, top, width, text, bg_color=None):
    """Add an insight/takeaway bar at the bottom of a slide."""
    if bg_color is None:
        bg_color = RGBColor(0xE0, 0xE7, 0xFF)  # Soft premium indigo/blue
    bar = _add_rounded_rect(slide, left, top, width, Inches(0.75), fill_color=bg_color)
    _add_textbox(slide, left + Inches(0.3), top + Inches(0.1),
                 width - Inches(0.6), Inches(0.55), text,
                 font_size=14, color=NAVY, alignment=PP_ALIGN.CENTER, italic=True)


def _add_slide_title(slide, title_text, subtitle_text=None):
    """Add the main title and optional subtitle to a slide."""
    _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.9),
                 title_text, font_size=28, bold=True, color=NAVY, font_name="Georgia")
    if subtitle_text:
        _add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.5),
                     subtitle_text, font_size=14, color=TEXT_MED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_1_title(prs, result):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide)

    # Decorative top line
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), fill_color=NAVY)

    # Main title
    _add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.4), Inches(2.2),
                 "Corporate Health Insurance\nReview: Utilization, Efficiency,\nand Risk Insights",
                 font_size=42, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER,
                 font_name="Georgia")

    # Subtitle
    _add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8.4), Inches(0.6),
                 "MIS Data Analysis & Executive Summary",
                 font_size=20, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

    # File name
    fname = result.get("file", "Unknown File")
    _add_textbox(slide, Inches(2.5), Inches(5.0), Inches(8.4), Inches(0.4),
                 f"Source: {fname}",
                 font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER, italic=True)

    # Bottom box
    _add_rounded_rect(slide, Inches(2), Inches(6.2), Inches(9.4), Inches(0.7),
                      fill_color=RGBColor(0xE8, 0xE5, 0xE5))
    _add_textbox(slide, Inches(2.3), Inches(6.3), Inches(8.8), Inches(0.5),
                 "Intended for Corporate Management, HR/Benefits, and Finance Leadership",
                 font_size=14, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)


def _slide_2_financial_overview(prs, result):
    """Slide 2: Financial Overview — KPIs + Cashless/Reimbursement split"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    kpi = result.get("kpis", {})

    total_claims = kpi.get("total_claims", 0)
    total_incurred = kpi.get("total_incurred", 0)
    cashless_count = kpi.get("cashless_count", 0)
    reimb_count = kpi.get("reimb_count", 0)
    cashless_pct = kpi.get("cashless_pct", 0)
    reimb_pct = round(100 - cashless_pct, 1)

    # Dynamic title
    title = f"Financial Overview: {cashless_pct}% cashless utilization"
    if cashless_pct > 60:
        title += " minimizes out-of-pocket employee stress"
    else:
        title += " — consider expanding network hospital coverage"
    _add_slide_title(slide, title)

    # Left side — Total Claims card
    _add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(3), Inches(2.3),
                      fill_color=WHITE)
    _add_textbox(slide, Inches(0.7), Inches(1.7), Inches(2.6), Inches(0.3),
                 "Total Claims", font_size=16, bold=True, color=NAVY)
    _add_textbox(slide, Inches(0.7), Inches(2.0), Inches(2.6), Inches(0.5),
                 str(total_claims), font_size=36, bold=True, color=NAVY)
    _add_textbox(slide, Inches(0.7), Inches(2.7), Inches(2.6), Inches(0.3),
                 "Total Approved Amount", font_size=14, bold=True, color=COPPER)
    _add_textbox(slide, Inches(0.7), Inches(3.0), Inches(2.6), Inches(0.5),
                 _fmt_inr_full(total_incurred), font_size=28, bold=True, color=NAVY)

    # Middle — Claim Volume breakdown
    _add_textbox(slide, Inches(4.2), Inches(1.6), Inches(4), Inches(0.4),
                 "Claim Volume", font_size=18, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    # Donut-like representation using shapes
    cx, cy = Inches(5.5), Inches(3.2)
    _add_shape_rect(slide, Inches(4.0), Inches(2.2), Inches(4.2), Inches(1.5),
                    fill_color=SLIDE_BG)

    # Cashless block
    _add_rounded_rect(slide, Inches(4.2), Inches(2.3), Inches(1.8), Inches(1.1),
                      fill_color=DARK_GREEN)
    _add_textbox(slide, Inches(4.3), Inches(2.4), Inches(1.6), Inches(0.9),
                 f"Cashless\n{cashless_count} Claims\n({cashless_pct}%)",
                 font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Reimbursement block
    _add_rounded_rect(slide, Inches(6.2), Inches(2.3), Inches(1.8), Inches(1.1),
                      fill_color=NAVY)
    _add_textbox(slide, Inches(6.3), Inches(2.4), Inches(1.6), Inches(0.9),
                 f"Reimbursement\n{reimb_count} Claims\n({reimb_pct}%)",
                 font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Right — Claim Value breakdown
    _add_textbox(slide, Inches(8.8), Inches(1.6), Inches(4), Inches(0.4),
                 "Claim Value", font_size=18, bold=True, color=NAVY, alignment=PP_ALIGN.CENTER)

    cashless_amt = total_incurred * (cashless_pct / 100) if cashless_pct > 0 else 0
    reimb_amt = total_incurred - cashless_amt

    _add_rounded_rect(slide, Inches(8.8), Inches(2.3), Inches(1.9), Inches(1.1),
                      fill_color=DARK_GREEN)
    _add_textbox(slide, Inches(8.9), Inches(2.4), Inches(1.7), Inches(0.9),
                 f"Cashless\n{_fmt_inr(cashless_amt)}\n({cashless_pct}%)",
                 font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_rounded_rect(slide, Inches(10.9), Inches(2.3), Inches(1.9), Inches(1.1),
                      fill_color=NAVY)
    _add_textbox(slide, Inches(11.0), Inches(2.4), Inches(1.7), Inches(0.9),
                 f"Reimbursement\n{_fmt_inr(reimb_amt)}\n({reimb_pct}%)",
                 font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Additional KPI row
    approval_rate = kpi.get("approval_rate", 0)
    avg_claim = kpi.get("avg_claim", 0)
    total_billed = kpi.get("total_billed", 0)
    deductions = kpi.get("total_deductions", 0)

    y = Inches(4.1)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    start_x = Inches(0.5)

    _add_kpi_card(slide, start_x, y, card_w, Inches(1.8),
                  "✅", f"{approval_rate}%", "APPROVAL RATE",
                  accent_color=DARK_GREEN,
                  sub_text=f"{kpi.get('approved_count', 0)} approved / {kpi.get('rejected_count', 0)} rejected")

    _add_kpi_card(slide, start_x + card_w + gap, y, card_w, Inches(1.8),
                  "📊", _fmt_inr(avg_claim), "AVG CLAIM SIZE",
                  accent_color=ACCENT_BLUE,
                  sub_text=f"Max: {_fmt_inr(kpi.get('max_claim', 0))}")

    _add_kpi_card(slide, start_x + 2 * (card_w + gap), y, card_w, Inches(1.8),
                  "💰", _fmt_inr(total_billed), "TOTAL BILLED",
                  accent_color=COPPER)

    _add_kpi_card(slide, start_x + 3 * (card_w + gap), y, card_w, Inches(1.8),
                  "📉", _fmt_inr(deductions), "TOTAL DEDUCTIONS",
                  accent_color=DANGER_RED)

    # Insight bar
    if cashless_pct > 60:
        insight = f"{cashless_pct}% of healthcare capital is deployed through cashless channels, demonstrating effective network utilization."
    else:
        insight = f"Only {cashless_pct}% cashless adoption. Expanding network hospital partnerships could reduce reimbursement burden."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_3_status_distribution(prs, result):
    """Slide 3: Claim Status Distribution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    status_dist = result.get("status_distribution", [])
    type_dist = result.get("claim_type_distribution", [])
    kpi = result.get("kpis", {})

    _add_slide_title(slide, "Claim Status & Type Distribution",
                     f"Overview of {kpi.get('total_claims', 0)} claims by approval status and claim type")

    # Left side — Status distribution table
    _add_textbox(slide, Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
                 "Status Distribution", font_size=16, bold=True, color=NAVY)

    if status_dist:
        headers = ["Status", "Count", "Total Amount"]
        rows = []
        for s in status_dist[:10]:
            rows.append([s.get("status", "—"), str(s.get("count", 0)),
                        _fmt_inr(s.get("total_amt", 0))])
        _add_table(slide, Inches(0.6), Inches(2.1), Inches(6), headers, rows)

    # Right side — Claim Type distribution
    _add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.4),
                 "Claim Type Split", font_size=16, bold=True, color=NAVY)

    if type_dist:
        headers = ["Type", "Count"]
        rows = [[t.get("type", "—"), str(t.get("count", 0))] for t in type_dist[:8]]
        _add_table(slide, Inches(7.2), Inches(2.1), Inches(4.5), headers, rows)

    # Insight
    approved = kpi.get("approved_count", 0)
    rejected = kpi.get("rejected_count", 0)
    approval_rate = kpi.get("approval_rate", 0)
    if approval_rate > 75:
        insight = f"Healthy approval rate of {approval_rate}%. {approved} claims approved, {rejected} rejected."
    else:
        insight = f"Approval rate at {approval_rate}% warrants review. {rejected} claims were rejected — consider analyzing rejection reasons."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_4_monthly_trend(prs, result):
    """Slide 4: Monthly Claims Trend"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    trend = result.get("monthly_trend", [])

    _add_slide_title(slide, "Monthly Claims Trend",
                     "Claim volume and cost patterns over time")

    if not trend:
        _add_textbox(slide, Inches(3), Inches(3), Inches(7), Inches(1),
                     "No date data available for trend analysis",
                     font_size=18, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
        return

    # Table with monthly data
    headers = ["Month", "Claims", "Total Amount", "Avg per Claim"]
    rows = []
    for t in trend:
        count = t.get("count", 0)
        total = t.get("total_amt", 0)
        avg = total / count if count > 0 else 0
        rows.append([t.get("month", "—"), str(count), _fmt_inr(total), _fmt_inr(avg)])

    _add_table(slide, Inches(0.6), Inches(1.8), Inches(10), headers, rows,
               col_widths=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])

    # Visual bars for each month (simplified bar chart)
    if len(trend) > 0:
        max_count = max(t.get("count", 1) for t in trend)
        bar_area_left = Inches(11.2)
        bar_area_top = Inches(1.8)
        bar_height = Inches(0.25)
        max_bar_width = Inches(1.8)
        spacing = Inches(0.35)

        _add_textbox(slide, bar_area_left, Inches(1.4), Inches(2), Inches(0.3),
                     "Volume", font_size=10, bold=True, color=TEXT_MED, alignment=PP_ALIGN.CENTER)

        for i, t in enumerate(trend):
            count = t.get("count", 0)
            bar_w = max(Inches(0.1), max_bar_width * (count / max_count)) if max_count > 0 else Inches(0.1)
            y_pos = bar_area_top + i * spacing
            _add_shape_rect(slide, bar_area_left, y_pos, bar_w, bar_height,
                           fill_color=NAVY)

    # Insight
    if len(trend) >= 2:
        first, last = trend[0], trend[-1]
        if last["count"] > first["count"]:
            pct = round((last["count"] - first["count"]) / first["count"] * 100, 1)
            insight = f"Claim volume grew {pct}% from {first['month']} to {last['month']}."
        else:
            pct = round((first["count"] - last["count"]) / first["count"] * 100, 1)
            insight = f"Claim volume declined {pct}% from {first['month']} to {last['month']}."
    else:
        insight = "Limited trend data available for analysis."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_5_hospital_breakdown(prs, result):
    """Slide 5: Top Hospitals (Facility Utilization)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    hosp = result.get("hospital_breakdown", [])
    kpi = result.get("kpis", {})

    top_hosp = hosp[0]["hospital"] if hosp else "N/A"
    _add_slide_title(slide, f"Facility Utilization: Top hospitals by claim volume",
                     f"Top {len(hosp)} hospitals ranked by number of claims")

    if not hosp:
        _add_textbox(slide, Inches(3), Inches(3), Inches(7), Inches(1),
                     "No hospital data available",
                     font_size=18, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
        return

    # Horizontal bar chart style — hospital names + bars
    max_count = max(h["count"] for h in hosp) if hosp else 1
    start_y = Inches(1.8)
    bar_spacing = Inches(0.42)
    name_width = Inches(4.5)
    bar_max_width = Inches(4.5)
    count_x = Inches(4.8)

    for i, h in enumerate(hosp[:10]):
        y = start_y + i * bar_spacing
        name = h.get("hospital", "Unknown")
        count = h.get("count", 0)
        amt = h.get("total_amt", 0)
        avg = amt / count if count > 0 else 0

        # Hospital name
        _add_textbox(slide, Inches(0.4), y, name_width, Inches(0.35),
                     name, font_size=10, color=TEXT_DARK, alignment=PP_ALIGN.RIGHT)

        # Bar
        bar_w = max(Inches(0.15), bar_max_width * (count / max_count))
        _add_shape_rect(slide, count_x, y + Inches(0.05), bar_w, Inches(0.25),
                       fill_color=NAVY)

        # Count label on bar
        _add_textbox(slide, count_x + bar_w + Inches(0.1), y, Inches(0.5), Inches(0.35),
                     str(count), font_size=10, bold=True, color=NAVY)

        # Amount on right
        _add_textbox(slide, Inches(10.5), y, Inches(2.5), Inches(0.35),
                     f"● {_fmt_inr(avg)} avg", font_size=9, color=COPPER)

    # Column headers
    _add_textbox(slide, count_x, Inches(1.45), Inches(2), Inches(0.3),
                 "Volume (No. of Claims)", font_size=10, bold=True, color=TEXT_MED)
    _add_textbox(slide, Inches(10.5), Inches(1.45), Inches(2.5), Inches(0.3),
                 "Average Cost (₹)", font_size=10, bold=True, color=TEXT_MED)

    # Insight
    if len(hosp) >= 2:
        h1, h2 = hosp[0], hosp[1]
        insight = (f"{h1['hospital']} leads with {h1['count']} claims ({_fmt_inr(h1['total_amt'])} total). "
                  f"Second highest: {h2['hospital']} with {h2['count']} claims.")
    else:
        insight = f"{top_hosp} is the primary provider."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_6_disease_breakdown(prs, result):
    """Slide 6: Disease / Diagnosis Breakdown"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    diseases = result.get("disease_breakdown", [])
    chronic = result.get("chronic_breakdown", [])

    _add_slide_title(slide, "Disease & Diagnosis Breakdown",
                     "Top disease categories by claim frequency and cost")

    # Left — Disease table
    if diseases:
        _add_textbox(slide, Inches(0.6), Inches(1.6), Inches(6), Inches(0.4),
                     "Top Disease Categories", font_size=14, bold=True, color=NAVY)
        headers = ["Category", "Claims", "Total Amount"]
        rows = []
        for d in diseases[:10]:
            rows.append([d.get("category", "—"), str(d.get("count", 0)),
                        _fmt_inr(d.get("total_amt", 0))])
        _add_table(slide, Inches(0.6), Inches(2.1), Inches(6.2), headers, rows)

    # Right — Chronic conditions
    if chronic:
        _add_textbox(slide, Inches(7.4), Inches(1.6), Inches(5.5), Inches(0.4),
                     "Chronic Conditions", font_size=14, bold=True, color=COPPER)
        headers = ["Condition", "Claims", "Cost"]
        rows = []
        for c in chronic[:8]:
            rows.append([c.get("illness", "—")[:35], str(c.get("count", 0)),
                        _fmt_inr(c.get("total_amt", 0))])
        _add_table(slide, Inches(7.4), Inches(2.1), Inches(5.2), headers, rows)

    # Insight
    if diseases:
        top = diseases[0]
        insight = f"Leading category: {top['category']} with {top['count']} claims totaling {_fmt_inr(top['total_amt'])}. Targeted wellness programs may reduce costs."
    else:
        insight = "No disease category data available for analysis."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_7_demographics(prs, result):
    """Slide 7: Demographics — Gender, Age, Relation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    gender = result.get("gender_breakdown", [])
    age = result.get("age_breakdown", [])
    relation = result.get("relation_distribution", [])
    kpi = result.get("kpis", {})

    _add_slide_title(slide, "Beneficiary Demographics",
                     "Gender, age group, and relationship distribution of claimants")

    col_x = Inches(0.5)
    card_w = Inches(3.8)

    # Gender section
    _add_textbox(slide, col_x, Inches(1.6), card_w, Inches(0.4),
                 "Gender Distribution", font_size=14, bold=True, color=NAVY)
    if gender:
        y = Inches(2.1)
        for g in gender:
            _add_rounded_rect(slide, col_x, y, Inches(3.5), Inches(0.5), fill_color=WHITE)
            _add_textbox(slide, col_x + Inches(0.2), y + Inches(0.08),
                         Inches(1.5), Inches(0.35), g["gender"],
                         font_size=12, color=TEXT_DARK)
            _add_textbox(slide, col_x + Inches(2.0), y + Inches(0.08),
                         Inches(1.3), Inches(0.35), str(g["count"]),
                         font_size=14, bold=True, color=NAVY, alignment=PP_ALIGN.RIGHT)
            y += Inches(0.6)

    # Age section
    _add_textbox(slide, col_x + Inches(4.3), Inches(1.6), card_w, Inches(0.4),
                 "Age Distribution", font_size=14, bold=True, color=NAVY)
    if age:
        headers = ["Age Group", "Count"]
        rows = [[a["group"], str(a["count"])] for a in age]
        _add_table(slide, col_x + Inches(4.3), Inches(2.1), Inches(3.5), headers, rows)

    # Relation section
    _add_textbox(slide, col_x + Inches(8.6), Inches(1.6), card_w, Inches(0.4),
                 "Beneficiary Relation", font_size=14, bold=True, color=NAVY)
    if relation:
        y = Inches(2.1)
        total = sum(r["count"] for r in relation)
        for r in relation:
            pct = round(r["count"] / total * 100, 1) if total > 0 else 0
            _add_rounded_rect(slide, col_x + Inches(8.6), y, Inches(3.5), Inches(0.5),
                             fill_color=WHITE)
            _add_textbox(slide, col_x + Inches(8.8), y + Inches(0.08),
                         Inches(1.5), Inches(0.35), r["relation"],
                         font_size=12, color=TEXT_DARK)
            _add_textbox(slide, col_x + Inches(10.5), y + Inches(0.08),
                         Inches(1.4), Inches(0.35), f"{r['count']} ({pct}%)",
                         font_size=12, bold=True, color=DARK_GREEN, alignment=PP_ALIGN.RIGHT)
            y += Inches(0.6)

    # Insight
    self_count = sum(r["count"] for r in relation if r.get("relation") == "Self")
    dep_count = sum(r["count"] for r in relation if r.get("relation") == "Dependents")
    total_claims = kpi.get("total_claims", 0)
    if self_count > 0 and dep_count > 0:
        insight = f"Employees (Self) account for {self_count} claims, Dependents for {dep_count} claims out of {total_claims} total."
    else:
        insight = f"Total pool: {total_claims} claims across all beneficiaries."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_8_fraud_flags(prs, result):
    """Slide 8: Fraud & Risk Flags"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    flags = result.get("fraud_flags", [])

    _add_slide_title(slide, f"Risk & Anomaly Flags: {len(flags)} claims flagged",
                     "Heuristic fraud/anomaly signals requiring investigation")

    if not flags:
        _add_rounded_rect(slide, Inches(3), Inches(3), Inches(7.4), Inches(1.2),
                         fill_color=RGBColor(0xD1, 0xFA, 0xE5))
        _add_textbox(slide, Inches(3.5), Inches(3.2), Inches(6.4), Inches(0.8),
                     "✅ No fraud/anomaly flags detected — All claims appear within normal parameters",
                     font_size=16, color=DARK_GREEN, alignment=PP_ALIGN.CENTER)
        return

    # Flags table
    headers = ["Claim ID", "Employee", "Hospital", "Amount", "Signals"]
    rows = []
    for f in flags[:10]:
        signals = "; ".join(f.get("signals", [])[:2])
        rows.append([
            str(f.get("claim_id", "—"))[:15],
            str(f.get("employee", "—"))[:20],
            str(f.get("hospital", "—"))[:20],
            _fmt_inr(f.get("amount", 0)),
            signals[:50]
        ])

    _add_table(slide, Inches(0.4), Inches(1.8), Inches(12.6), headers, rows,
               col_widths=[Inches(1.8), Inches(2.5), Inches(2.5), Inches(1.8), Inches(4)])

    insight = f"{len(flags)} claims flagged for anomalies. These require further investigation before final settlement."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight,
                     bg_color=RGBColor(0xFE, 0xE2, 0xE2))


def _slide_9_high_value(prs, result):
    """Slide 9: High-Value Claims"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    high_vals = result.get("high_value_claims", [])
    outliers = result.get("outliers", [])

    _add_slide_title(slide, f"High-Value & Outlier Claims",
                     f"{len(high_vals)} high-value claims, {len(outliers)} statistical outliers detected")

    # High-value table
    if high_vals:
        _add_textbox(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
                     "Highest-Value Claims", font_size=14, bold=True, color=COPPER)
        headers = ["Claim ID", "Employee", "Hospital", "Amount"]
        rows = []
        for h in high_vals[:10]:
            rows.append([
                str(h.get("claim_id", "—"))[:15],
                str(h.get("employee", "—"))[:22],
                str(h.get("hospital", "—"))[:22],
                _fmt_inr(h.get("amount", 0))
            ])
        _add_table(slide, Inches(0.5), Inches(2.1), Inches(6.2), headers, rows)

    # Outliers table
    if outliers:
        _add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.4),
                     "Statistical Outliers (Z > 2.5)", font_size=14, bold=True, color=DANGER_RED)
        headers = ["Employee", "Hospital", "Amount", "Z-Score"]
        rows = []
        for o in outliers[:8]:
            rows.append([
                str(o.get("employee", "—"))[:18],
                str(o.get("hospital", "—"))[:18],
                _fmt_inr(o.get("amount", 0)),
                str(o.get("z_score", "—"))
            ])
        _add_table(slide, Inches(7.2), Inches(2.1), Inches(5.4), headers, rows)

    # Insight
    total_hv = sum(h.get("amount", 0) for h in high_vals)
    insight = f"{len(high_vals)} high-value claims totaling {_fmt_inr(total_hv)}. Monitoring these cases can help manage claim reserves."
    _add_insight_bar(slide, Inches(0.5), Inches(6.3), Inches(12.4), insight)


def _slide_10_conclusions(prs, result):
    """Slide 10: Strategic Conclusions & AI Narrative"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    kpi = result.get("kpis", {})
    narrative = result.get("ai_narrative", "")

    _add_slide_title(slide, "Strategic Conclusions",
                     "Key findings and recommendations")

    # Three conclusion cards
    cashless_pct = kpi.get("cashless_pct", 0)
    approval_rate = kpi.get("approval_rate", 0)
    diseases = result.get("disease_breakdown", [])
    top_disease = diseases[0]["category"] if diseases else "N/A"

    cards = [
        {
            "icon": "⚙️",
            "title": "Operational Resilience",
            "text": (f"Cashless adoption at {cashless_pct}%. "
                    f"Approval rate: {approval_rate}%. "
                    f"Review rejection patterns to improve claimant experience.")
        },
        {
            "icon": "💰",
            "title": "Financial Utilization",
            "text": (f"Total incurred: {_fmt_inr(kpi.get('total_incurred', 0))}. "
                    f"Average claim: {_fmt_inr(kpi.get('avg_claim', 0))}. "
                    f"Deductions total: {_fmt_inr(kpi.get('total_deductions', 0))}.")
        },
        {
            "icon": "🏥",
            "title": "Clinical Intervention",
            "text": (f"Leading disease category: {top_disease}. "
                    f"Proactive corporate wellness interventions are recommended "
                    f"to curb future premium inflation.")
        },
    ]

    card_w = Inches(12.4)
    card_h = Inches(1.0)
    start_y = Inches(1.8)

    for i, c in enumerate(cards):
        y = start_y + i * (card_h + Inches(0.2))
        _add_rounded_rect(slide, Inches(0.5), y, card_w, card_h, fill_color=WHITE)

        # Icon + Title
        _add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.5), Inches(0.4),
                     c["icon"], font_size=20)
        _add_textbox(slide, Inches(1.5), y + Inches(0.15), Inches(4), Inches(0.35),
                     c["title"], font_size=16, bold=True, color=NAVY)

        # Body text
        _add_textbox(slide, Inches(1.5), y + Inches(0.45), Inches(10.8), Inches(0.55),
                     c["text"], font_size=12, color=TEXT_DARK)

    # AI narrative at bottom (trimmed)
    if narrative:
        clean = narrative.replace("**", "").replace("*", "")[:300]
        _add_textbox(slide, Inches(0.5), Inches(5.6), Inches(12.4), Inches(1.2),
                     clean, font_size=10, color=TEXT_MED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ═══════════════════════════════════════════════════════════════════════════════

def generate_ppt(result: dict) -> io.BytesIO:
    """
    Generate an editable PowerPoint presentation from analysis results.
    Returns a BytesIO buffer containing the .pptx file.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all 10 slides
    _slide_1_title(prs, result)
    _slide_2_financial_overview(prs, result)
    _slide_3_status_distribution(prs, result)
    _slide_4_monthly_trend(prs, result)
    _slide_5_hospital_breakdown(prs, result)
    _slide_6_disease_breakdown(prs, result)
    _slide_7_demographics(prs, result)
    _slide_8_fraud_flags(prs, result)
    _slide_9_high_value(prs, result)
    _slide_10_conclusions(prs, result)

    # Save to buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
